from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Circle
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / 'results'
METRICS = RESULTS / 'metrics.csv'
OUT_PPT = RESULTS / 'DATS5990_Poster_Licheng_Guo.pptx'

sns.set_theme(style='whitegrid', context='talk')

NAVY = RGBColor(15, 32, 62)
BLUE = RGBColor(32, 99, 155)
CYAN = RGBColor(77, 166, 255)
SLATE = RGBColor(81, 92, 111)
LIGHT = RGBColor(246, 249, 252)
PHASE1 = RGBColor(245, 166, 35)
PHASE2 = RGBColor(46, 204, 113)
PHASE3 = RGBColor(142, 68, 173)
RED = RGBColor(220, 76, 100)
WHITE = RGBColor(255, 255, 255)

DISPLAY_NAMES = {
    'popularity': 'Popularity',
    'user_cf': 'User CF',
    'item_cf': 'Item CF',
    'matrix_factorization': 'MF (SVD)',
    'ncf': 'NCF',
    'hybrid': 'Hybrid',
}


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_panel(slide, x, y, w, h, title, color=NAVY, fill=WHITE, title_size=22, radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = color
    shape.line.width = Pt(1.4)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.10), w - Inches(0.36), Inches(0.35), title, size=title_size, bold=True, color=color)
    return shape


def add_bullets(slide, x, y, w, h, title, bullets, color=NAVY, fill=WHITE, font_size=17):
    add_panel(slide, x, y, w, h, title, color=color, fill=fill)
    box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.55), w - Inches(0.30), h - Inches(0.65))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = NAVY
        p.bullet = True


def draw_box(ax, xy, w, h, text, fc, ec='#1f2d3d', fs=12, weight='bold'):
    patch = FancyBboxPatch(xy, w, h, boxstyle='round,pad=0.02,rounding_size=0.03', facecolor=fc, edgecolor=ec, linewidth=2)
    ax.add_patch(patch)
    ax.text(xy[0] + w / 2, xy[1] + h / 2, text, ha='center', va='center', fontsize=fs, fontweight=weight)


def arrow(ax, start, end, color='#34495e'):
    ax.add_patch(FancyArrowPatch(start, end, arrowstyle='-|>', mutation_scale=16, linewidth=2, color=color))


def make_architecture_figures():
    RESULTS.mkdir(parents=True, exist_ok=True)

    # Phase pipeline
    fig, ax = plt.subplots(figsize=(12, 3.6))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    draw_box(ax, (0.03, 0.22), 0.25, 0.56, 'Phase I\nBaseline Models\nPopularity / User-CF / Item-CF / MF', '#FDEBD0', fs=14)
    draw_box(ax, (0.37, 0.22), 0.23, 0.56, 'Phase II\nNeural Modeling\nNCF with user/item embeddings + MLP', '#D5F5E3', fs=14)
    draw_box(ax, (0.69, 0.22), 0.28, 0.56, 'Phase III\nHybrid + Fairness\nContent fusion + long-tail reranking', '#EBDEF0', fs=14)
    arrow(ax, (0.28, 0.50), (0.37, 0.50))
    arrow(ax, (0.60, 0.50), (0.69, 0.50))
    ax.text(0.5, 0.93, 'Three-Stage Recommender Research Pipeline', ha='center', fontsize=18, fontweight='bold')
    plt.tight_layout()
    plt.savefig(RESULTS / 'pipeline_diagram.png', dpi=240, bbox_inches='tight')
    plt.close()

    # NCF architecture
    fig, ax = plt.subplots(figsize=(10, 4.6))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    draw_box(ax, (0.03, 0.32), 0.16, 0.22, 'User ID', '#D6EAF8', fs=13)
    draw_box(ax, (0.03, 0.62), 0.16, 0.22, 'Item ID', '#D6EAF8', fs=13)
    draw_box(ax, (0.26, 0.28), 0.18, 0.27, 'User\nEmbedding', '#AED6F1', fs=13)
    draw_box(ax, (0.26, 0.58), 0.18, 0.27, 'Item\nEmbedding', '#AED6F1', fs=13)
    draw_box(ax, (0.51, 0.41), 0.15, 0.28, 'Concatenate', '#D4EFDF', fs=13)
    draw_box(ax, (0.72, 0.28), 0.12, 0.18, 'MLP\n128', '#F9E79F', fs=12)
    draw_box(ax, (0.72, 0.51), 0.12, 0.18, 'MLP\n64', '#F9E79F', fs=12)
    draw_box(ax, (0.89, 0.39), 0.08, 0.20, 'Score', '#F5B7B1', fs=13)
    arrow(ax, (0.19, 0.43), (0.26, 0.43))
    arrow(ax, (0.19, 0.73), (0.26, 0.73))
    arrow(ax, (0.44, 0.43), (0.51, 0.50))
    arrow(ax, (0.44, 0.72), (0.51, 0.59))
    arrow(ax, (0.66, 0.55), (0.72, 0.60))
    arrow(ax, (0.78, 0.51), (0.78, 0.46))
    arrow(ax, (0.84, 0.50), (0.89, 0.50))
    ax.text(0.5, 0.95, 'Neural Collaborative Filtering Architecture', ha='center', fontsize=18, fontweight='bold')
    plt.tight_layout()
    plt.savefig(RESULTS / 'ncf_architecture.png', dpi=240, bbox_inches='tight')
    plt.close()

    # Hybrid architecture
    fig, ax = plt.subplots(figsize=(10, 4.8))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    draw_box(ax, (0.03, 0.58), 0.18, 0.18, 'User–Item\nInteractions', '#D6EAF8', fs=13)
    draw_box(ax, (0.03, 0.22), 0.18, 0.18, 'Product\nDescriptions', '#D6EAF8', fs=13)
    draw_box(ax, (0.28, 0.56), 0.20, 0.22, 'Collaborative\nSignal\n(MF/SVD)', '#FDEBD0', fs=13)
    draw_box(ax, (0.28, 0.20), 0.20, 0.22, 'Content\nSignal\n(TF-IDF)', '#D5F5E3', fs=13)
    draw_box(ax, (0.58, 0.37), 0.16, 0.22, 'Fusion\nWeighted Sum', '#EBDEF0', fs=13)
    draw_box(ax, (0.81, 0.37), 0.15, 0.22, 'Fairness /\nDiversity\nReranking', '#F5B7B1', fs=13)
    arrow(ax, (0.21, 0.67), (0.28, 0.67))
    arrow(ax, (0.21, 0.31), (0.28, 0.31))
    arrow(ax, (0.48, 0.67), (0.58, 0.52))
    arrow(ax, (0.48, 0.31), (0.58, 0.44))
    arrow(ax, (0.74, 0.48), (0.81, 0.48))
    ax.text(0.5, 0.95, 'Hybrid Recommender and Re-ranking Architecture', ha='center', fontsize=18, fontweight='bold')
    plt.tight_layout()
    plt.savefig(RESULTS / 'hybrid_architecture.png', dpi=240, bbox_inches='tight')
    plt.close()


def make_charts(df: pd.DataFrame):
    df = df.copy()
    df['display'] = df['model'].map(DISPLAY_NAMES)

    plt.figure(figsize=(11, 5.6))
    plot_df = df.melt(id_vars=['display'], value_vars=['precision@k', 'recall@k', 'ndcg@k', 'map@k'], var_name='Metric', value_name='Score')
    sns.barplot(data=plot_df, x='display', y='Score', hue='Metric', palette='Blues')
    plt.title('Accuracy Comparison on Online Retail E-commerce Dataset', fontsize=18, fontweight='bold')
    plt.xlabel('')
    plt.ylabel('Metric value', fontsize=13)
    plt.xticks(rotation=15, fontsize=12)
    plt.yticks(fontsize=12)
    plt.legend(title='')
    plt.tight_layout()
    plt.savefig(RESULTS / 'accuracy_metrics.png', dpi=240)
    plt.close()

    plt.figure(figsize=(11, 5.6))
    trade_df = df.melt(id_vars=['display'], value_vars=['coverage', 'diversity', 'fairness_gap'], var_name='Metric', value_name='Score')
    sns.barplot(data=trade_df, x='display', y='Score', hue='Metric', palette=['#2E86C1', '#48C9B0', '#AF7AC5'])
    plt.title('Coverage, Diversity, and Fairness Trade-offs', fontsize=18, fontweight='bold')
    plt.xlabel('')
    plt.ylabel('Metric value', fontsize=13)
    plt.xticks(rotation=15, fontsize=12)
    plt.yticks(fontsize=12)
    plt.legend(title='')
    plt.tight_layout()
    plt.savefig(RESULTS / 'tradeoff_metrics.png', dpi=240)
    plt.close()

    ranking = df.sort_values('ndcg@k', ascending=False)
    plt.figure(figsize=(8, 5.8))
    sns.barplot(data=ranking, y='display', x='ndcg@k', hue='display', dodge=False, palette='viridis', legend=False)
    plt.title('Model Ranking by NDCG@10', fontsize=18, fontweight='bold')
    plt.xlabel('NDCG@10', fontsize=13)
    plt.ylabel('')
    plt.xticks(fontsize=12)
    plt.yticks(fontsize=12)
    plt.tight_layout()
    plt.savefig(RESULTS / 'ranking_ndcg.png', dpi=240)
    plt.close()


def build_poster(df: pd.DataFrame):
    prs = Presentation()
    prs.slide_width = Inches(24)
    prs.slide_height = Inches(36)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = WHITE

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.55))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.color.rgb = NAVY

    add_textbox(slide, Inches(0.65), Inches(0.32), Inches(16.7), Inches(0.85), 'A Systematic Study of Recommendation Algorithms for E-commerce Platforms', size=30, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.68), Inches(1.22), Inches(11), Inches(0.5), 'Licheng Guo  ·  DATS5990 Independent Study  ·  University of Pennsylvania', size=16, color=WHITE)
    add_textbox(slide, Inches(15.7), Inches(0.34), Inches(7.6), Inches(1.5), 'Real e-commerce dataset + three-stage recommender benchmark + academic poster redesign', size=16, color=WHITE, align=PP_ALIGN.RIGHT)

    add_bullets(slide, Inches(0.6), Inches(2.9), Inches(7.4), Inches(2.95), 'Research Motivation', [
        'E-commerce platforms require recommenders that balance ranking quality, product coverage, diversity, and fairness.',
        'Purely accuracy-oriented methods often overexpose head products and under-serve the long tail.',
        'This project studies whether hybrid and fairness-aware methods improve practical recommendation quality.'
    ], color=BLUE, fill=RGBColor(248, 251, 255), font_size=18)

    add_bullets(slide, Inches(8.25), Inches(2.9), Inches(7.4), Inches(2.95), 'Dataset and Protocol', [
        'Dataset: UCI Online Retail, a transactional e-commerce dataset with customers, stock codes, timestamps, and product descriptions.',
        'Users with at least 5 interactions and items with at least 5 purchases were retained.',
        'Evaluation uses chronological train/validation/test splitting and Top-K ranking metrics.'
    ], color=BLUE, fill=RGBColor(248, 251, 255), font_size=18)

    add_bullets(slide, Inches(15.9), Inches(2.9), Inches(7.45), Inches(2.95), 'Main Contributions', [
        'Unified pipeline for baseline, neural, and hybrid recommenders.',
        'Switched from movie benchmark to a true product-level e-commerce dataset.',
        'Poster integrates model architectures, empirical results, and deployment-oriented interpretation.'
    ], color=RED, fill=RGBColor(255, 249, 250), font_size=18)

    add_panel(slide, Inches(0.6), Inches(6.15), Inches(22.75), Inches(3.1), 'Overall Research Pipeline', color=BLUE, fill=WHITE, title_size=24)
    slide.shapes.add_picture(str(RESULTS / 'pipeline_diagram.png'), Inches(1.0), Inches(6.8), width=Inches(21.9))

    add_panel(slide, Inches(0.6), Inches(9.55), Inches(7.25), Inches(6.3), 'Phase I — Classical Baselines', color=PHASE1, fill=RGBColor(255, 251, 245), title_size=23)
    add_bullets(slide, Inches(0.82), Inches(10.15), Inches(6.8), Inches(5.45), 'Included Models', [
        'Popularity ranking',
        'User-based collaborative filtering',
        'Item-based collaborative filtering',
        'Matrix factorization via truncated SVD',
        'Role: establish interpretable and competitive baselines for offline ranking.'
    ], color=PHASE1, fill=RGBColor(255, 251, 245), font_size=18)

    add_panel(slide, Inches(8.1), Inches(9.55), Inches(7.25), Inches(6.3), 'Phase II — Neural Modeling', color=PHASE2, fill=RGBColor(245, 255, 247), title_size=23)
    slide.shapes.add_picture(str(RESULTS / 'ncf_architecture.png'), Inches(8.35), Inches(10.15), width=Inches(6.8))
    add_textbox(slide, Inches(8.45), Inches(14.15), Inches(6.55), Inches(1.15), 'NCF learns nonlinear user–item interactions using embedding layers followed by an MLP scoring network.', size=17, color=SLATE)

    add_panel(slide, Inches(15.6), Inches(9.55), Inches(7.75), Inches(6.3), 'Phase III — Hybrid + Fairness-aware Optimization', color=PHASE3, fill=RGBColor(251, 247, 255), title_size=23)
    slide.shapes.add_picture(str(RESULTS / 'hybrid_architecture.png'), Inches(15.9), Inches(10.12), width=Inches(7.1))
    add_textbox(slide, Inches(16.0), Inches(14.15), Inches(6.95), Inches(1.15), 'Hybrid ranking combines collaborative and content signals, then applies reranking to improve long-tail exposure and recommendation diversity.', size=17, color=SLATE)

    add_panel(slide, Inches(0.6), Inches(16.2), Inches(11.15), Inches(6.15), 'Quantitative Results: Accuracy Metrics', color=BLUE, fill=WHITE, title_size=24)
    slide.shapes.add_picture(str(RESULTS / 'accuracy_metrics.png'), Inches(0.9), Inches(16.95), width=Inches(10.55))

    add_panel(slide, Inches(12.05), Inches(16.2), Inches(11.3), Inches(6.15), 'Quantitative Results: Trade-off Metrics', color=BLUE, fill=WHITE, title_size=24)
    slide.shapes.add_picture(str(RESULTS / 'tradeoff_metrics.png'), Inches(12.35), Inches(16.95), width=Inches(10.7))

    add_panel(slide, Inches(0.6), Inches(22.65), Inches(8.4), Inches(5.2), 'Ranking Summary', color=BLUE, fill=WHITE, title_size=24)
    slide.shapes.add_picture(str(RESULTS / 'ranking_ndcg.png'), Inches(0.95), Inches(23.35), width=Inches(7.6))

    best_acc = df.sort_values('ndcg@k', ascending=False).iloc[0]
    best_cov = df.sort_values('coverage', ascending=False).iloc[0]
    best_div = df.sort_values('diversity', ascending=False).iloc[0]

    add_bullets(slide, Inches(9.3), Inches(22.65), Inches(7.0), Inches(5.2), 'Key Empirical Findings', [
        f'Best ranking accuracy: {DISPLAY_NAMES[best_acc.model]} (NDCG@10 = {best_acc["ndcg@k"]:.3f}).',
        f'Best coverage: {DISPLAY_NAMES[best_cov.model]} (Coverage = {best_cov.coverage:.3f}).',
        f'Best diversity: {DISPLAY_NAMES[best_div.model]} (Diversity = {best_div.diversity:.3f}).',
        'Hybrid recommendation offers the strongest overall balance across ranking quality, catalog coverage, and diversity.'
    ], color=RED, fill=RGBColor(255, 249, 250), font_size=18)

    add_bullets(slide, Inches(16.55), Inches(22.65), Inches(6.8), Inches(5.2), 'Interpretation for E-commerce', [
        'Neighborhood and factorization methods remain strong on sparse transactional data.',
        'Neural models are feasible, but are not automatically superior without richer features or sequence signals.',
        'Hybrid reranking is attractive for deployment because it broadens product exposure while preserving accuracy.'
    ], color=PHASE3, fill=RGBColor(251, 247, 255), font_size=18)

    add_panel(slide, Inches(0.6), Inches(28.15), Inches(22.75), Inches(4.15), 'Metric Table', color=BLUE, fill=WHITE, title_size=24)
    rows, cols = len(df) + 1, 8
    table = slide.shapes.add_table(rows, cols, Inches(0.9), Inches(28.95), Inches(22.15), Inches(2.95)).table
    headers = ['Model', 'Precision@10', 'Recall@10', 'NDCG@10', 'MAP@10', 'Coverage', 'Diversity', 'Fairness Gap']
    widths = [2.8, 2.35, 2.2, 2.15, 2.1, 2.1, 2.1, 2.45]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.color.rgb = WHITE
    show = df.sort_values('ndcg@k', ascending=False)
    for i, row in enumerate(show.itertuples(index=False), start=1):
        vals = [DISPLAY_NAMES[row.model], row._0, row._1, row._2, row._3, row.coverage, row.diversity, row.fairness_gap]
        for j, val in enumerate(vals):
            cell = table.cell(i, j)
            cell.text = str(val) if isinstance(val, str) else f'{val:.3f}'
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 251, 255)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
                    r.font.color.rgb = NAVY

    add_bullets(slide, Inches(0.6), Inches(32.6), Inches(10.9), Inches(2.55), 'Conclusion', [
        'Using a real e-commerce dataset changes the framing from generic recommendation to product-level retail recommendation.',
        'The hybrid model achieved the best overall balance, especially for catalog coverage and diversity.',
        'This supports the argument that practical recommender systems should optimize more than accuracy alone.'
    ], color=BLUE, fill=RGBColor(248, 251, 255), font_size=18)

    add_bullets(slide, Inches(11.8), Inches(32.6), Inches(11.55), Inches(2.55), 'Future Directions', [
        'Add sequential models such as GRU4Rec or Transformer recommenders.',
        'Integrate richer product attributes and user history features.',
        'Extend fairness analysis to seller exposure, category fairness, and cold-start product promotion.'
    ], color=PHASE3, fill=RGBColor(251, 247, 255), font_size=18)

    add_textbox(slide, Inches(0.7), Inches(35.35), Inches(22.0), Inches(0.35), 'Poster regenerated from the updated Online Retail e-commerce experiment pipeline. Layout intentionally styled closer to academic conference posters with architecture diagrams and enlarged typography.', size=11, color=SLATE, align=PP_ALIGN.CENTER)

    prs.save(OUT_PPT)
    print(f'Saved poster to {OUT_PPT}')


def main():
    df = pd.read_csv(METRICS)
    make_architecture_figures()
    make_charts(df)
    build_poster(df)


if __name__ == '__main__':
    main()
